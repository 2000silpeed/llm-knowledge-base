"""PowerPoint 인제스터 (W1-04) — 멀티모달 2-패스 버전

텍스트 패스: python-pptx로 슬라이드별 텍스트/테이블/이미지 캡션 추출
이미지 패스: LibreOffice headless → PyMuPDF로 슬라이드 PNG 렌더링 → Vision LLM 상세 분석
조립:       슬라이드별 텍스트 + ### 시각 분석 섹션 병합 → 최종 마크다운

필수 외부 도구 (이미지 패스 활성화 시):
  - LibreOffice (libreoffice 명령)
  - PyMuPDF (fitz, pyproject.toml 의존성에 포함)

설정:
  ingest.slide_render: true         # 슬라이드 이미지 렌더링 활성화 (기본: true)
  vision_llm:                        # (선택) 별도 Vision 모델 설정, 없으면 기본 llm 사용
    provider: ollama
    model: gemma3:4b
    base_url: http://localhost:11434

사용 예:
    from scripts.ingest_ppt import ingest_ppt
    result = ingest_ppt("/path/to/file.pptx")
    # {"status": "ok", "path": "raw/office/file.md", "slides": 24, ...}
"""

import hashlib
import logging
import re
import subprocess
import tempfile
import unicodedata
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from pathlib import Path

import yaml

from scripts.token_counter import estimate_tokens, load_settings, parse_frontmatter

logger = logging.getLogger(__name__)


def _slugify(text: str, max_len: int = 60) -> str:
    """텍스트를 파일명용 슬러그로 변환합니다."""
    text = unicodedata.normalize("NFKD", text)
    text = re.sub(r"[^\w\s가-힣]", "", text, flags=re.UNICODE)
    text = re.sub(r"\s+", "-", text.strip())
    text = text.lower()
    text = re.sub(r"-+", "-", text).strip("-")
    return text[:max_len] if text else "document"


# ──────────────────────────────────────────────
# Vision 설정 헬퍼
# ──────────────────────────────────────────────

def _get_vision_settings(settings: dict) -> dict:
    """비전 분석용 settings 반환.

    settings에 'vision_llm' 키가 있으면 해당 설정으로 llm을 대체합니다.
    없으면 기본 llm 설정을 그대로 사용합니다.
    """
    if "vision_llm" in settings:
        return {**settings, "llm": settings["vision_llm"]}
    return settings


# ──────────────────────────────────────────────
# 슬라이드 이미지 렌더링 (LibreOffice + PyMuPDF)
# ──────────────────────────────────────────────

def _render_slides_to_bytes(pptx_path: Path, work_dir: Path) -> list[bytes]:
    """PPTX → PDF → 슬라이드별 PNG bytes 변환 (디스크 PNG 저장 없음).

    LibreOffice headless로 PPTX를 PDF로 변환한 뒤,
    PyMuPDF(fitz)로 페이지별 PNG bytes를 메모리에서 생성합니다.

    Args:
        pptx_path: 원본 .pptx 파일 경로
        work_dir:  LibreOffice PDF 출력 디렉토리 (임시 디렉토리)

    Returns:
        슬라이드 순서대로 정렬된 PNG bytes 목록

    Raises:
        RuntimeError: LibreOffice 또는 PyMuPDF 오류
    """
    import fitz  # PyMuPDF

    pdf_path = work_dir / (pptx_path.stem + ".pdf")

    try:
        result = subprocess.run(
            [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", str(work_dir), str(pptx_path),
            ],
            capture_output=True,
            timeout=120,
        )
    except FileNotFoundError:
        raise RuntimeError(
            "LibreOffice를 찾을 수 없습니다. "
            "'sudo apt install libreoffice' 또는 'brew install libreoffice'로 설치하세요."
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError("LibreOffice 변환 타임아웃 (120초 초과)")

    if result.returncode != 0:
        err = result.stderr.decode(errors="replace").strip()
        raise RuntimeError(f"LibreOffice 변환 실패 (code={result.returncode}): {err}")

    if not pdf_path.exists():
        raise RuntimeError(f"LibreOffice 변환 후 PDF를 찾을 수 없습니다: {pdf_path}")

    # PDF → 슬라이드별 PNG bytes (PyMuPDF, 2x 해상도, 디스크 저장 없음)
    mat = fitz.Matrix(2.0, 2.0)  # 144 DPI — 차트/텍스트 선명도 확보
    with fitz.open(str(pdf_path)) as doc:
        slide_bytes = [
            doc[page_num].get_pixmap(matrix=mat).tobytes("png")
            for page_num in range(len(doc))
        ]

    pdf_path.unlink(missing_ok=True)
    logger.info("슬라이드 렌더링 완료: %d장", len(slide_bytes))
    return slide_bytes


# ──────────────────────────────────────────────
# Vision 캡션 (기존 — 임베드 이미지용)
# ──────────────────────────────────────────────

def _generate_caption(image_bytes: bytes, ext: str, settings: dict) -> str:
    """Vision API로 슬라이드 임베드 이미지의 한 문장 캡션을 생성합니다."""
    try:
        from scripts.llm import call_vision

        media_map = {
            "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "png": "image/png", "gif": "image/gif", "webp": "image/webp",
            "bmp": "image/png",
        }
        media_type = media_map.get(ext.lower(), "image/png")
        prompt = (
            "이 프레젠테이션 슬라이드를 한 문장으로 간결하게 설명해주세요. "
            "차트·그래프·표가 있으면 핵심 수치와 결론을 포함하세요."
        )
        return call_vision(image_bytes, media_type, prompt, settings).strip()
    except Exception as exc:
        logger.warning("Vision 캡션 생성 실패: %s", exc)
        return ""


# ──────────────────────────────────────────────
# Vision 슬라이드 상세 분석 (신규)
# ──────────────────────────────────────────────

def _analyze_slide_image(image_bytes: bytes, slide_num: int, settings: dict) -> str:
    """슬라이드 전체 이미지를 Vision LLM으로 상세 분석합니다.

    텍스트 추출이 놓친 시각적 정보(차트·다이어그램·레이아웃·강조)를 보완합니다.

    Args:
        image_bytes: 슬라이드 PNG 바이너리
        slide_num:   슬라이드 번호 (1-based, 로그용)
        settings:    load_settings() 결과 (vision_llm 있으면 자동 사용)

    Returns:
        마크다운 형식의 분석 텍스트. 오류 시 빈 문자열.
    """
    from scripts.llm import call_vision

    vision_settings = _get_vision_settings(settings)
    prompt = (
        "이 프레젠테이션 슬라이드를 분석하여 마크다운으로 정리하세요.\n\n"
        "다음 항목을 순서대로 작성하세요 (해당 내용이 있는 항목만):\n"
        "- **주제**: 슬라이드의 핵심 메시지 (1~2문장)\n"
        "- **텍스트**: 슬라이드에 표시된 주요 텍스트·레이블·수치 (목록으로)\n"
        "- **차트/그래프**: 종류, 축 레이블, 주요 데이터 포인트, 추세\n"
        "- **표**: 행/열 구조와 주요 셀 값\n"
        "- **다이어그램/흐름도**: 구성 요소와 상호 관계\n"
        "- **시각적 강조**: 색상·크기·배치로 부각된 정보\n\n"
        "없는 항목은 생략하고 한국어로 답하세요."
    )

    try:
        result = call_vision(image_bytes, "image/png", prompt, vision_settings).strip()
        logger.debug("슬라이드 %d 시각 분석 완료 (%d자)", slide_num, len(result))
        return result
    except Exception as exc:
        logger.warning("슬라이드 %d 시각 분석 실패 (건너뜀): %s", slide_num, exc)
        return ""


def _run_vision_pass(slide_bytes_list: list[bytes], settings: dict) -> list[str]:
    """슬라이드 이미지 배열을 병렬로 Vision LLM 분석합니다.

    concept_extractor.py의 ThreadPoolExecutor 패턴과 동일.
    Ollama는 GPU 요청을 서버 측에서 직렬화하지만,
    I/O·HTTP 오버헤드를 overlap하여 전체 시간을 단축합니다.
    """
    analyses: list[str] = [""] * len(slide_bytes_list)

    def _analyze(args: tuple[int, bytes]) -> tuple[int, str]:
        idx, img_bytes = args
        return idx, _analyze_slide_image(img_bytes, idx + 1, settings)

    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = {
            executor.submit(_analyze, (i, b)): i
            for i, b in enumerate(slide_bytes_list)
        }
        for future in as_completed(futures):
            try:
                idx, analysis = future.result()
                analyses[idx] = analysis
            except Exception as e:
                slide_num = futures[future] + 1
                logger.warning("슬라이드 %d Vision 분석 실패 (건너뜀): %s", slide_num, e)

    return analyses


# ──────────────────────────────────────────────
# 슬라이드 → 마크다운 (텍스트 패스)
# ──────────────────────────────────────────────

def _shape_to_text(shape) -> str:
    """도형에서 텍스트를 추출합니다."""
    if not shape.has_text_frame:
        return ""

    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        level = para.level
        font_size = None
        for run in para.runs:
            if run.font and run.font.size:
                try:
                    font_size = run.font.size.pt
                except Exception:
                    pass
                break

        if font_size and font_size >= 28:
            lines.append(f"### {text}")
        elif font_size and font_size >= 20:
            lines.append(f"#### {text}")
        elif level == 0:
            lines.append(text)
        else:
            indent = "  " * level
            lines.append(f"{indent}- {text}")

    return "\n".join(lines)


def _table_shape_to_markdown(shape) -> str:
    """테이블 도형을 마크다운 테이블로 변환합니다."""
    if not shape.has_table:
        return ""

    table = shape.table
    rows_md: list[str] = []

    for row_idx, row in enumerate(table.rows):
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("|", "\\|").replace("\n", " ")
            cells.append(cell_text)
        row_line = "| " + " | ".join(cells) + " |"
        rows_md.append(row_line)
        if row_idx == 0:
            separator = "| " + " | ".join(["---"] * len(cells)) + " |"
            rows_md.append(separator)

    return "\n".join(rows_md)


def _slide_to_markdown(
    slide,
    slide_num: int,
    images_dir: Path,
    do_vision: bool,
    settings: dict,
) -> tuple[str, list[str]]:
    """슬라이드 텍스트 패스: python-pptx로 텍스트/테이블/임베드 이미지 추출.

    Returns:
        (markdown_text, list_of_saved_image_paths)
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    sections: list[str] = []
    saved_images: list[str] = []

    # 제목 추출
    title_text = ""
    body_shapes = []

    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
            if ph is not None and ph.idx in (0, 1):
                title_text = shape.text.strip()
                continue
        except Exception:
            pass
        body_shapes.append(shape)

    heading = f"## Slide {slide_num}" + (f": {title_text}" if title_text else "")
    sections.append(heading)

    # 본문 도형 처리
    for shape in body_shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_blob = shape.image.blob
                img_ext = shape.image.ext or "png"
                img_hash = hashlib.md5(img_blob).hexdigest()[:12]
                filename = f"{img_hash}.{img_ext}"
                dest = images_dir / filename
                if not dest.exists():
                    dest.write_bytes(img_blob)
                    logger.debug("임베드 이미지 저장: slide%d → %s", slide_num, dest)
                saved_images.append(str(dest))

                caption = ""
                if do_vision:
                    caption = _generate_caption(img_blob, img_ext, settings)
                sections.append(f"![{caption or '슬라이드 이미지'}](raw/images/{filename})")
                continue

            if shape.has_table:
                table_md = _table_shape_to_markdown(shape)
                if table_md:
                    sections.append(table_md)
                continue

            if shape.has_text_frame:
                text = _shape_to_text(shape)
                if text:
                    sections.append(text)

        except Exception as exc:
            logger.debug("도형 처리 건너뜀 (slide%d): %s", slide_num, exc)

    # 발표자 노트
    try:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip() if notes_slide else ""
        if notes_text:
            note_lines = "\n".join(f"> {line}" for line in notes_text.splitlines())
            sections.append(f"\n> **Note:**\n{note_lines}")
    except Exception:
        pass

    return "\n\n".join(filter(None, sections)), saved_images


# ──────────────────────────────────────────────
# 목차 생성
# ──────────────────────────────────────────────

def _build_toc(slide_titles: list[tuple[int, str]]) -> str:
    """슬라이드 번호+제목 목록으로 목차를 생성합니다."""
    lines = ["## 목차\n"]
    for num, title in slide_titles:
        lines.append(f"- Slide {num}" + (f": {title}" if title else ""))
    return "\n".join(lines)


def _extract_slide_title(slide) -> str:
    """슬라이드에서 제목 텍스트를 추출합니다."""
    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
            if ph is not None and ph.idx in (0, 1):
                return shape.text.strip()
        except Exception:
            pass
    return ""


# ──────────────────────────────────────────────
# 슬라이드 조립 (텍스트 + 시각 분석)
# ──────────────────────────────────────────────

def _assemble_slide(text_md: str, visual_analysis: str) -> str:
    """텍스트 패스 결과와 시각 분석 결과를 하나의 슬라이드 섹션으로 조립합니다."""
    if not visual_analysis:
        return text_md

    return f"{text_md}\n\n### 시각 분석\n\n{visual_analysis}"


# ──────────────────────────────────────────────
# 퍼블릭 진입점
# ──────────────────────────────────────────────

def ingest_ppt(
    ppt_path: str | Path,
    project_root: Path | str | None = None,
    settings: dict | None = None,
) -> dict:
    """PowerPoint 파일(.pptx)을 인제스트합니다 — 2-패스 멀티모달 버전.

    텍스트 패스: python-pptx 기반 텍스트·테이블·임베드 이미지 추출
    이미지 패스: LibreOffice 렌더링 + Vision LLM 슬라이드 상세 분석 (선택적)
    조립:       슬라이드별 병합 → 최종 마크다운 파일 생성

    Args:
        ppt_path:     PowerPoint 파일 경로
        project_root: 프로젝트 루트 경로. None이면 이 파일 기준 상위 디렉토리.
        settings:     설정 dict. None이면 settings.yaml 자동 로드.

    Returns:
        {
            "status":       "ok" | "error",
            "path":         str,          # 저장된 .md 파일 경로
            "meta_path":    str,          # 저장된 .meta.yaml 파일 경로
            "title":        str,
            "slides":       int,
            "token_count":  int,
            "images":       list[str],
            "visual_pass":  bool,         # 이미지 패스 성공 여부
            "message":      str,          # 오류 시 메시지
        }
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"status": "error", "message": "python-pptx가 설치되지 않았습니다. `uv add python-pptx`"}

    if project_root is None:
        project_root = Path(__file__).parent.parent
    project_root = Path(project_root)
    ppt_path = Path(ppt_path)

    if not ppt_path.exists():
        return {"status": "error", "message": f"파일 없음: {ppt_path}"}
    if ppt_path.suffix.lower() not in {".pptx", ".ppt"}:
        return {"status": "error", "message": f"PowerPoint 파일이 아닙니다: {ppt_path}"}

    if settings is None:
        settings = load_settings()

    office_dir = project_root / settings["paths"]["raw"] / "office"
    images_dir = project_root / settings["paths"]["raw"] / "images"
    office_dir.mkdir(parents=True, exist_ok=True)
    images_dir.mkdir(parents=True, exist_ok=True)

    ingest_cfg = settings.get("ingest", {})
    chunking_cfg = settings.get("chunking", {})

    slides_per_chunk = chunking_cfg.get("ppt_slides_per_chunk", 10)
    do_vision = ingest_cfg.get("vision_caption", True)
    do_slide_render = ingest_cfg.get("slide_render", True)

    logger.info("PowerPoint 로드 중: %s", ppt_path)

    try:
        prs = Presentation(str(ppt_path))
    except Exception as exc:
        return {"status": "error", "message": f"PowerPoint 열기 실패: {exc}"}

    slides = prs.slides
    slide_count = len(slides)
    collected_at = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

    # 프레젠테이션 제목
    title = ppt_path.stem
    try:
        core_props = prs.core_properties
        if core_props.title:
            title = core_props.title.strip() or title
    except Exception:
        pass

    # 슬라이드 제목 목록 (목차용)
    slide_titles: list[tuple[int, str]] = []
    for i, slide in enumerate(slides, start=1):
        slide_titles.append((i, _extract_slide_title(slide)))

    toc_text = _build_toc(slide_titles)

    # ── 텍스트 패스 ──
    logger.info("텍스트 패스 시작 (%d슬라이드)", slide_count)
    text_mds: list[str] = []
    all_images: list[str] = []

    for i, slide in enumerate(slides, start=1):
        logger.debug("텍스트 패스 슬라이드 처리: %d/%d", i, slide_count)
        slide_md, slide_images = _slide_to_markdown(
            slide, i, images_dir, do_vision, settings
        )
        text_mds.append(slide_md)
        all_images.extend(slide_images)

    # ── 이미지 패스 (선택적) ──
    visual_analyses: list[str] = [""] * slide_count
    visual_pass_ok = False

    if do_slide_render:
        logger.info("이미지 패스 시작 — LibreOffice 렌더링 후 Vision 분석")
        try:
            with tempfile.TemporaryDirectory(prefix="kb_ppt_") as tmp_dir:
                slide_bytes_list = _render_slides_to_bytes(ppt_path, Path(tmp_dir))
                visual_analyses = _run_vision_pass(slide_bytes_list, settings)
            visual_pass_ok = True
            logger.info("이미지 패스 완료")
        except Exception as exc:
            logger.warning("이미지 패스 실패 — 텍스트 패스만 사용합니다: %s", exc)
    else:
        logger.info("이미지 패스 비활성화 (ingest.slide_render: false)")

    # ── 조립 패스 ──
    assembled_mds = [_assemble_slide(t, v) for t, v in zip(text_mds, visual_analyses)]

    # ── 청크 분할 + 목차 반복 ──
    chunk_count = max(1, -(-slide_count // slides_per_chunk))
    chunks: list[str] = []

    for chunk_idx in range(chunk_count):
        start = chunk_idx * slides_per_chunk
        end = min(start + slides_per_chunk, slide_count)
        chunk_slides = assembled_mds[start:end]

        if chunk_count > 1:
            chunk_header = (
                f"<!-- chunk {chunk_idx + 1}/{chunk_count}: "
                f"Slide {start + 1}–{end} -->\n\n"
                f"{toc_text}\n\n"
                f"---\n"
            )
        else:
            chunk_header = f"{toc_text}\n\n---\n"

        chunk_body = "\n\n---\n\n".join(chunk_slides)
        chunks.append(chunk_header + "\n" + chunk_body)

    body = "\n\n---\n\n".join(chunks)
    token_count = estimate_tokens(body)

    # 이미지 경로 정규화
    rel_images: list[str] = []
    for img_path in all_images:
        try:
            rel_images.append(str(Path(img_path).relative_to(project_root)))
        except ValueError:
            rel_images.append(img_path)

    frontmatter: dict = {
        "title": title,
        "source_file": str(ppt_path),
        "collected_at": collected_at,
        "slides": slide_count,
        "token_count": token_count,
        "images": rel_images,
        "visual_pass": visual_pass_ok,
    }
    fm_str = yaml.dump(frontmatter, allow_unicode=True, default_flow_style=False, sort_keys=False)
    document = f"---\n{fm_str}---\n\n# {title}\n\n{body}\n"

    # 파일명 결정
    stem = _slugify(ppt_path.stem) or "presentation"
    filename = f"{stem}.md"
    dest = office_dir / filename
    if dest.exists():
        file_hash = hashlib.md5(ppt_path.stem.encode()).hexdigest()[:6]
        filename = f"{stem}_{file_hash}.md"
        dest = office_dir / filename

    dest.write_text(document, encoding="utf-8")
    rel_path = str(dest.relative_to(project_root))

    # .meta.yaml 저장
    meta: dict = {
        "source_file": str(ppt_path),
        "collected_at": collected_at,
        "title": title,
        "slide_count": slide_count,
        "slides_per_chunk": slides_per_chunk,
        "chunk_count": chunk_count,
        "token_count": token_count,
        "images": rel_images,
        "visual_pass": visual_pass_ok,
        "slide_titles": [{"num": n, "title": t} for n, t in slide_titles],
    }
    meta_filename = dest.stem + ".meta.yaml"
    meta_dest = office_dir / meta_filename
    meta_dest.write_text(
        yaml.dump(meta, allow_unicode=True, default_flow_style=False, sort_keys=False),
        encoding="utf-8",
    )
    meta_rel_path = str(meta_dest.relative_to(project_root))

    logger.info(
        "저장 완료: %s (슬라이드: %d, 토큰: %d, 시각 분석: %s)",
        rel_path, slide_count, token_count, "완료" if visual_pass_ok else "미실행",
    )

    return {
        "status": "ok",
        "path": rel_path,
        "meta_path": meta_rel_path,
        "title": title,
        "slides": slide_count,
        "token_count": token_count,
        "images": rel_images,
        "visual_pass": visual_pass_ok,
    }


# ──────────────────────────────────────────────
# Vision 재실행 (W1-04c)
# ──────────────────────────────────────────────

def _parse_ppt_md(text: str) -> tuple[dict, str]:
    """인제스트된 PPT .md 파일의 frontmatter와 body를 파싱합니다."""
    return parse_frontmatter(text)


def _find_slides_without_vision(body: str) -> list[int]:
    """body에서 시각 분석이 없거나 비어있는 슬라이드 번호(1-based) 목록 반환."""
    slide_matches = list(re.finditer(r"^## Slide (\d+)", body, re.MULTILINE))
    missing: list[int] = []

    for i, match in enumerate(slide_matches):
        sec_start = match.start()
        sec_end = slide_matches[i + 1].start() if i + 1 < len(slide_matches) else len(body)
        section = body[sec_start:sec_end]

        # ### 시각 분석 존재 + 내용 있는지 확인
        visual_match = re.search(r"### 시각 분석\s*\n\n(\S)", section)
        if not visual_match:
            missing.append(int(match.group(1)))

    return missing


def _inject_visual_analysis(body: str, slide_num: int, analysis: str) -> str:
    """body의 특정 슬라이드 섹션에 시각 분석을 주입하거나 교체합니다.

    - 기존 `### 시각 분석` 섹션이 있으면 교체
    - 없으면 슬라이드 섹션 끝(---구분자 앞)에 추가
    """
    all_slides = list(re.finditer(r"^## Slide (\d+)", body, re.MULTILINE))
    target_idx = next(
        (i for i, m in enumerate(all_slides) if int(m.group(1)) == slide_num),
        None,
    )
    if target_idx is None:
        return body

    sec_start = all_slides[target_idx].start()
    sec_end = all_slides[target_idx + 1].start() if target_idx + 1 < len(all_slides) else len(body)
    section = body[sec_start:sec_end]

    new_visual = f"\n\n### 시각 분석\n\n{analysis.strip()}"

    existing = re.search(r"\n\n### 시각 분석\n\n", section)
    if existing:
        # 기존 분석 끝 위치 탐색: 다음 헤더 또는 --- 구분자
        after_start = existing.end()
        after_text = section[after_start:]
        next_boundary = re.search(r"\n\n(?:#{1,4} |---)", after_text)
        if next_boundary:
            cut = after_start + next_boundary.start()
            new_section = section[:existing.start()] + new_visual + section[cut:]
        else:
            # 분석이 섹션 끝까지 이어짐 — trailing 구분자 보존
            trail = re.search(r"(\n\n---\s*)$", section)
            if trail:
                new_section = section[:existing.start()] + new_visual + section[trail.start():]
            else:
                new_section = section[:existing.start()] + new_visual
    else:
        # 없으면 섹션 끝 --- 앞에 삽입
        trail = re.search(r"(\n\n---\s*)$", section)
        if trail:
            new_section = section[:trail.start()] + new_visual + section[trail.start():]
        else:
            new_section = section.rstrip() + new_visual

    return body[:sec_start] + new_section + body[sec_end:]


def retry_vision_pass(
    md_path: str | Path,
    *,
    pptx_path: str | Path | None = None,
    project_root: Path | str | None = None,
    settings: dict | None = None,
    force: bool = False,
    only_slides: list[int] | None = None,
    dry_run: bool = False,
) -> dict:
    """이미 인제스트된 PPT .md 파일에 Vision 이미지 분석을 재실행합니다.

    기본 동작 (force=False):
      - `visual_pass: true` 파일은 건너뜀
      - `### 시각 분석`이 없는 슬라이드만 선택적으로 재실행

    Args:
        md_path:     재실행할 raw/office/{name}.md 파일 경로
        pptx_path:   원본 PPTX 경로. None이면 frontmatter source_file 자동 참조.
        project_root: 프로젝트 루트. None이면 스크립트 상위 디렉토리.
        settings:    설정 dict. None이면 settings.yaml 자동 로드.
        force:       visual_pass: true여도 강제 재실행
        only_slides: 재실행할 슬라이드 번호 목록 (1-based). None이면 자동 탐지.
        dry_run:     실제 실행 없이 대상 슬라이드 목록만 반환.

    Returns:
        {
            "status":          "ok" | "error",
            "md_path":         str,
            "pptx_path":       str,
            "target_slides":   list[int],   # 재실행 대상 슬라이드
            "slides_done":     int,         # 성공한 슬라이드 수
            "slides_failed":   int,         # Vision 실패 슬라이드 수
            "visual_pass":     bool,        # 완료 후 전체 분석 여부
            "message":         str,
        }
    """
    if project_root is None:
        project_root = Path(__file__).parent.parent
    project_root = Path(project_root)
    md_path = Path(md_path)

    if not md_path.exists():
        return {"status": "error", "message": f"파일 없음: {md_path}"}

    if settings is None:
        settings = load_settings()

    # frontmatter 파싱
    text = md_path.read_text(encoding="utf-8")
    meta, body = _parse_ppt_md(text)

    # visual_pass 체크
    if meta.get("visual_pass") and not force:
        return {
            "status": "ok",
            "md_path": str(md_path),
            "pptx_path": str(pptx_path or meta.get("source_file", "")),
            "target_slides": [],
            "slides_done": 0,
            "slides_failed": 0,
            "visual_pass": True,
            "message": "이미 시각 분석이 완료된 파일입니다. --force로 강제 재실행하세요.",
        }

    # 원본 PPTX 경로 확인
    if pptx_path is None:
        pptx_path = Path(meta.get("source_file", ""))
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        return {
            "status": "error",
            "message": (
                f"원본 PPTX 파일을 찾을 수 없습니다: {pptx_path}\n"
                "--pptx 옵션으로 경로를 직접 지정하세요."
            ),
        }

    # 재실행 대상 슬라이드 결정
    if only_slides:
        target_slides = sorted(set(only_slides))
    elif force:
        all_slide_nums = [
            int(m.group(1))
            for m in re.finditer(r"^## Slide (\d+)", body, re.MULTILINE)
        ]
        target_slides = all_slide_nums
    else:
        target_slides = _find_slides_without_vision(body)

    result_base: dict = {
        "status": "ok",
        "md_path": str(md_path),
        "pptx_path": str(pptx_path),
        "target_slides": target_slides,
        "slides_done": 0,
        "slides_failed": 0,
        "visual_pass": bool(meta.get("visual_pass")),
        "message": "",
    }

    if not target_slides:
        result_base["message"] = "재실행할 슬라이드가 없습니다 (모든 슬라이드에 시각 분석 존재)."
        return result_base

    if dry_run:
        result_base["message"] = f"재실행 예정 슬라이드: {target_slides}"
        return result_base

    # 슬라이드 렌더링 (전체 → 필요한 장만 Vision 호출)
    logger.info("슬라이드 렌더링 시작: %s", pptx_path)
    try:
        with tempfile.TemporaryDirectory(prefix="kb_ppt_retry_") as tmp_dir:
            all_slide_bytes = _render_slides_to_bytes(pptx_path, Path(tmp_dir))
    except Exception as exc:
        return {**result_base, "status": "error", "message": f"슬라이드 렌더링 실패: {exc}"}

    logger.info("Vision 분석 시작: 슬라이드 %s", target_slides)
    slides_done = 0
    slides_failed = 0

    # 범위 초과 슬라이드 사전 필터링
    valid_slides = []
    for slide_num in target_slides:
        idx = slide_num - 1
        if idx < 0 or idx >= len(all_slide_bytes):
            logger.warning("슬라이드 %d: 범위 초과 (총 %d장) — 건너뜀", slide_num, len(all_slide_bytes))
            slides_failed += 1
        else:
            valid_slides.append(slide_num)

    # Vision 분석 병렬 실행
    vision_settings = _get_vision_settings(settings)
    analyses: dict[int, str] = {}

    def _analyze_one(slide_num: int) -> tuple[int, str]:
        return slide_num, _analyze_slide_image(all_slide_bytes[slide_num - 1], slide_num, vision_settings)

    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = {executor.submit(_analyze_one, sn): sn for sn in valid_slides}
        for future in as_completed(futures):
            try:
                slide_num, analysis = future.result()
                analyses[slide_num] = analysis
            except Exception as exc:
                slide_num = futures[future]
                logger.warning("슬라이드 %d Vision 분석 예외: %s", slide_num, exc)
                analyses[slide_num] = ""

    # 결과 주입 (순서 보장을 위해 정렬)
    for slide_num in sorted(analyses):
        analysis = analyses[slide_num]
        if analysis:
            body = _inject_visual_analysis(body, slide_num, analysis)
            slides_done += 1
            logger.debug("슬라이드 %d 시각 분석 주입 완료", slide_num)
        else:
            slides_failed += 1
            logger.warning("슬라이드 %d 시각 분석 빈 결과 — 건너뜀", slide_num)

    # visual_pass 상태 재계산
    remaining_missing = _find_slides_without_vision(body)
    all_done = len(remaining_missing) == 0

    # frontmatter 갱신
    meta["visual_pass"] = all_done
    fm_str = yaml.dump(meta, allow_unicode=True, default_flow_style=False, sort_keys=False)
    md_path.write_text(f"---\n{fm_str}---\n\n{body}", encoding="utf-8")

    # .meta.yaml 갱신
    meta_path = md_path.with_suffix(".meta.yaml")
    if meta_path.exists():
        try:
            meta_yaml = yaml.safe_load(meta_path.read_text(encoding="utf-8")) or {}
            meta_yaml["visual_pass"] = all_done
            meta_path.write_text(
                yaml.dump(meta_yaml, allow_unicode=True, default_flow_style=False, sort_keys=False),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.warning("meta.yaml 갱신 실패: %s", exc)

    logger.info(
        "Vision 재실행 완료 — 성공: %d, 실패: %d, visual_pass: %s",
        slides_done, slides_failed, all_done,
    )

    return {
        **result_base,
        "slides_done": slides_done,
        "slides_failed": slides_failed,
        "visual_pass": all_done,
        "message": "" if all_done else f"미완료 슬라이드: {remaining_missing}",
    }
